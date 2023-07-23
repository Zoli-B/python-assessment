import model
from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import os
import logging


logging.basicConfig(filename='presentation.log', level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')


prs = Presentation()


for item in model.json_data['presentation']:
    slide_type = item['type']

    logging.info(f"Processing slide of type: {slide_type}")

    #Title slide
    if slide_type == 'title':
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = item['title']
        subtitle.text = item['content']

        logging.info(f"Title slide - Title: {item['title']}")
        logging.info(f"Title slide - Content: {item['content']}")


    #Text_slide
    if slide_type == 'text':
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = item['title']
        content.text = item['content']

        logging.info(f"Text slide - Title: {item['title']}")
        logging.info(f"Text slide - Content: {item['content']}")

    #List_slide
    if slide_type == 'list':
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = item['title']
        frame = content.text_frame

        for element in item['content']:
            p = frame.add_paragraph()
            p.text = element['text']
            p.level = element['level']

        logging.info(f"List slide - Title: {item['title']}")
        logging.info("List slide - List elements:")
        for element in item['content']:
            logging.info(f"    - {element['text']} (level: {element['level']})")

    #Picture slide
    if slide_type == 'picture':
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = item['title']
        image_name = item['content']
        image_path = os.path.join(os.getcwd(), image_name)
        left = Inches(2)
        top = Inches(1.5)
        width = Inches(6)
        height = Inches(4.5)
        slide.shapes.add_picture(image_path, left, top, width,height)

        logging.info(f"Picture slide - Title: {item['title']}")
        logging.info(f"Picture slide - Image path: {image_path}")

    #Plot slide
    if slide_type == 'plot':
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = item['title']
        x = []
        y = []

        for row in model.csv_data:
            x.append(float(row[0]))
            y.append(float(row[1]))

        fig, ax = plt.subplots()
        ax.plot(x,y)

        for label in model.json_data["presentation"]:
            if 'configuration' in label:
                config = label['configuration']
                x = config.get("x-label")
                y = config.get("y-label")


        ax.set_xlabel(x)
        ax.set_ylabel(y)

        plot_path = "plot.png"
        plt.savefig(plot_path)

        left = Inches(1)
        top = Inches(1.5)
        width = Inches(7)
        height = Inches(5)

        logging.info(f"Plot slide - Title: {item['title']}")
        logging.info(f"Plot slide - X Label: {x}")
        logging.info(f"Plot slide - Y Label: {y}")

        #lábjegyzet
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = "Felhasznált forrás: Google, ChatGPT, Stackoverflow"

        slide.shapes.add_picture(plot_path, left, top, width, height)
